# extractor.py (sketch)
from pathlib import Path
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_path

def extract_docx(path):
    doc = docx.Document(path)
    blocks = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if text:
            # could inspect p.style.name for headings
            blocks.append({"type":"paragraph", "text": text})
    return blocks

def extract_pptx(path):
    prs = Presentation(path)
    blocks = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        if slide_text:
            blocks.append({"type":"slide", "index": i, "text": "\n".join(slide_text)})
    return blocks

def extract_pdf(path):
    blocks = []
    with pdfplumber.open(path) as pdf:
        for pageno, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            if text.strip():
                blocks.append({"type":"page", "pageno": pageno, "text": text})
            else:
                # fallback to OCR for this page
                images = convert_from_path(path, first_page=pageno+1, last_page=pageno+1)
                # usually one image
                text = pytesseract.image_to_string(images[0])
                if text.strip():
                    blocks.append({"type":"page_ocr", "pageno": pageno, "text": text})
    return blocks

def extract_image(path):
    return [{"type":"image", "text": pytesseract.image_to_string(Image.open(path))}]
